"""Extract reviewable text from office/binary docs without entering PPT production.

Usage:
  python scripts/extract-office-text.py INPUT OUTPUT.md
"""

from __future__ import annotations

import argparse
import re
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree

DEFAULT_MAX_FILE_BYTES = 100 * 1024 * 1024
SENSITIVE_PATH_PATTERNS = (
    re.compile(r"(?i)(^|/)\.env(\..+)?$"),
    re.compile(r"(?i)(^|/)(credentials?|secrets?|tokens?|passwords?)(\.[^/]*)?$"),
    re.compile(r"(?i)(^|/)(service[-_]?account|gcp[-_]?credentials?)([-_.][^/]*)?\.json$"),
    re.compile(r"(?i)(^|/)(id_rsa|id_dsa|id_ecdsa|id_ed25519)(\.pub)?$"),
    re.compile(r"(?i)\.(pem|p12|pfx|key|jks|keystore)$"),
)
TAG_RE = re.compile(r"<[^>]+>")


def is_sensitive_path(path: Path) -> bool:
    normalized = path.as_posix()
    return any(pattern.search(normalized) for pattern in SENSITIVE_PATH_PATTERNS)


def extract_docx(path: Path) -> list[str]:
    from docx import Document

    doc = Document(path)
    lines: list[str] = [f"# Extracted DOCX: {path.name}", ""]
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)
    for index, table in enumerate(doc.tables, 1):
        lines.extend(["", f"## Table {index}", ""])
        for row in table.rows:
            cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
            lines.append(" | ".join(cells))
    return lines


def extract_pptx(path: Path) -> list[str]:
    from pptx import Presentation

    deck = Presentation(path)
    lines: list[str] = [f"# Extracted PPTX evidence: {path.name}", ""]
    for slide_number, slide in enumerate(deck.slides, 1):
        lines.extend([f"## Slide {slide_number}", ""])
        found = False
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if text:
                lines.append(text)
                found = True
        if slide.has_notes_slide:
            notes = []
            for shape in slide.notes_slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    notes.append(shape.text.strip())
            if notes:
                lines.extend(["", "**Notes**", *notes])
                found = True
        if not found:
            lines.append("_No extractable text; visual review remains required._")
        lines.append("")
    return lines


def extract_xlsx(path: Path) -> list[str]:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    lines: list[str] = [f"# Extracted XLSX: {path.name}", ""]
    for sheet in workbook.worksheets:
        lines.extend([f"## {sheet.title}", ""])
        found = False
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if cell is None else str(cell).replace("\n", " ").strip() for cell in row]
            if any(cells):
                lines.append(" | ".join(cells))
                found = True
        if not found:
            lines.append("_No extractable cells; visual review remains required._")
        lines.append("")
    return lines


def extract_pdf(path: Path) -> list[str]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    lines: list[str] = [f"# Extracted PDF: {path.name}", ""]
    if len(reader.pages) == 0:
        lines.append("_No pages; visual review remains required._")
        return lines
    for index, page in enumerate(reader.pages, 1):
        text = (page.extract_text() or "").strip()
        lines.extend(
            [f"## Page {index}", "", text or "_No extractable text; visual review remains required._", ""]
        )
    return lines


def extract_drawio(path: Path) -> list[str]:
    raw = path.read_text(encoding="utf-8", errors="replace")
    text = TAG_RE.sub(" ", raw)
    text = re.sub(r"\s+", " ", text).strip()
    lines = [f"# Extracted draw.io XML labels: {path.name}", ""]
    if text:
        lines.append(text)
    else:
        lines.append("_No extractable labels; visual review remains required._")
    return lines


def extract_vsdx(path: Path) -> list[str]:
    lines: list[str] = [f"# Extracted VSDX XML text: {path.name}", ""]
    found = False
    with zipfile.ZipFile(path) as archive:
        names = [name for name in archive.namelist() if name.lower().endswith(".xml")]
        for name in sorted(names):
            xml_text = archive.read(name).decode("utf-8", errors="replace")
            labels = []
            try:
                root = ElementTree.fromstring(xml_text)
                for node in root.iter():
                    if node.text and node.text.strip():
                        labels.append(node.text.strip())
            except ElementTree.ParseError:
                stripped = TAG_RE.sub(" ", xml_text)
                labels = [part for part in stripped.split() if part]
            if labels:
                found = True
                lines.extend([f"## {name}", "", " ".join(labels[:400]), ""])
    if not found:
        lines.append("_No extractable XML text; convert to PNG/PDF and re-inventory._")
    return lines


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Extract evidence text from office/binary docs without entering PPT production."
    )
    parser.add_argument("input")
    parser.add_argument("output")
    parser.add_argument(
        "--max-file-bytes",
        type=int,
        default=DEFAULT_MAX_FILE_BYTES,
        help=f"refuse inputs larger than this (default: {DEFAULT_MAX_FILE_BYTES})",
    )
    args = parser.parse_args()
    if args.max_file_bytes < 1:
        parser.error("--max-file-bytes must be positive")

    source = Path(args.input).resolve()
    output = Path(args.output).resolve()
    if not source.is_file():
        print(f"input is not a file: {source}", file=sys.stderr)
        return 2
    if is_sensitive_path(source):
        print(
            "refusing potential secret/auth/private-key material; input was not opened",
            file=sys.stderr,
        )
        return 3
    size = source.stat().st_size
    if size > args.max_file_bytes:
        print(
            f"unsupported: input bytes {size} exceed MaxFileBytes "
            f"{args.max_file_bytes}; provide a reviewed extract or raise the limit",
            file=sys.stderr,
        )
        return 3
    suffix = source.suffix.lower()

    if suffix == ".ppt":
        print(
            "legacy .ppt requires human conversion to .pptx or PDF, then re-inventory; "
            "do not enter the PPT production stage",
            file=sys.stderr,
        )
        return 4

    extractors = {
        ".docx": extract_docx,
        ".pptx": extract_pptx,
        ".xlsx": extract_xlsx,
        ".pdf": extract_pdf,
        ".drawio": extract_drawio,
        ".vsdx": extract_vsdx,
    }
    extractor = extractors.get(suffix)
    if extractor is None:
        print(f"unsupported Office format: {suffix}", file=sys.stderr)
        return 2

    lines = extractor(source)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text("\n".join(lines), encoding="utf-8")
    print(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
